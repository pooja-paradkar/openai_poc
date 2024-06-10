import logging
from pptx.util import Pt
from flask import Flask, request, send_file, make_response, jsonify
from openai import AzureOpenAI
from pptx import Presentation
from app import api_key, azure_endpoint_url, deployment_model_GPT



app = Flask(__name__)

logging.basicConfig(filename='error.log', level=logging.ERROR)

client_ppt = AzureOpenAI(
    azure_endpoint=azure_endpoint_url,
    api_key=api_key,
    api_version="2024-02-15-preview"
)

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify(status='healthy')

def generate_content(topic):
    # Call OpenAI API to generate content based on the provided topic
    response = client_ppt.chat.completions.create(
        model=deployment_model_GPT,
        messages=[
            {
                "content": topic,
                'role': 'user'
            },
        ]
    )
    content = response.choices[0].message.content
    return content

def create_presentation(topic):
    # Generate content based on the topic provided
    content = generate_content(topic)

    # Create a new presentation
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = topic

    # Add content slides
    bullet_slide_layout = prs.slide_layouts[1]
    content_text = content.split('\n')
    for point in content_text:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Content"
        
        # Add text to the content slide with bullet points
        tf = shapes.placeholders[1].text_frame
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)  # Adjust font size as needed
        p.space_after = Pt(10)  # Adjust spacing as needed
        p.level = 0  # Set the bullet point level

    # Save the presentation to a file
    filename = f"{topic}_presentation.pptx"
    prs.save(filename)
    return filename

# @app.route('/download_ppt', methods=['GET'])
# def download_presentation():
#     data = request.get_json()
#     topic = data.get('topic')

#     if not topic:
#         return 'Topic not provided', 400  # Bad request if topic is missing

#     # Generate the presentation
#     filename = create_presentation(topic)
 
#     # Create an HTTP response
#     response = make_response(send_file(filename))
#     response.headers['Content-Type'] = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
#     response.headers['Content-Disposition'] = f'attachment; filename={filename}'
 
#     return response

@app.route('/download_ppt', methods=['GET'])
def download_presentation():
    data = request.get_json()
    topic = data.get('topic')

    if not topic:
        return 'Topic not provided', 400  # Bad request if the topic is missing

    # Generate the presentation
    filename = create_presentation(topic)

    # Create an HTTP response
    response = make_response(send_file(filename))
    response.headers['Content-Type'] = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    response.headers['Content-Disposition'] = f'attachment; filename={filename}'

    return response


if __name__ == '__main__':
    logging.basicConfig(level=logging.INFO)
    app.run(port=5000,debug=True)